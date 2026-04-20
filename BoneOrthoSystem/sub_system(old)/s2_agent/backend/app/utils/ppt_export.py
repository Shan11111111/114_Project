
from io import BytesIO
from typing import Dict, List

from pptx import Presentation
from pptx.util import Pt


TITLE_FONT_SIZE = Pt(32)
BODY_FONT_SIZE = Pt(20)


def _set_title_font(slide) -> None:
    if not slide.shapes.title:
        return
    tf = slide.shapes.title.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = TITLE_FONT_SIZE


def _fill_bullets(text_frame, lines: List[str]) -> None:
    # 清空原本 placeholder 的文字
    text_frame.clear()
    first = True
    for line in lines:
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        p.text = line
        p.level = 0
        for run in p.runs:
            run.font.size = BODY_FONT_SIZE


def create_pptx_from_summary(session_id: str, summary: Dict[str, List[str]]) -> BytesIO:
    """
    根據 llm.summarize_for_report 的結果產生 PPTX。
    summary = {
      "questions": [...],
      "assistant_points": [...],
      "health_tips": [...],
    }
    """
    prs = Presentation()

    # ---- 投影片 1：封面 ----
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide.shapes.title.text = "骨科互動助理 - 學習報告"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Session ID：{session_id}"

    _set_title_font(slide)
    # 副標題就一起用 BODY_FONT_SIZE
    for p in subtitle.text_frame.paragraphs:
        for run in p.runs:
            run.font.size = BODY_FONT_SIZE

    # ---- 投影片 2：本次提問概要 ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = "本次提問概要"
    _set_title_font(slide)

    body = slide.placeholders[1].text_frame
    q_lines = summary.get("questions", []) or ["（本次提問未偵測到明確問題）"]
    _fill_bullets(body, q_lines)

    # ---- 投影片 3：助理說明重點 ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "助理說明重點"
    _set_title_font(slide)

    body = slide.placeholders[1].text_frame
    a_lines = summary.get("assistant_points", []) or ["（目前沒有可摘要的說明內容）"]
    _fill_bullets(body, a_lines)

    # ---- 投影片 4：一般健康建議 ----
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "一般健康說明與建議"
    _set_title_font(slide)

    body = slide.placeholders[1].text_frame
    h_lines = summary.get("health_tips", []) or [
        "若有實際不適症狀，仍建議儘速就醫，由專業骨科或復健科醫師評估。",
    ]
    _fill_bullets(body, h_lines)

    # ---- 輸出成 BytesIO ----
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio