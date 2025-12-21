# BoneOrthoBackend/s2_agent/legacy_agent/backend/app/routers/export.py
from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import List

from fastapi import APIRouter
from fastapi.responses import StreamingResponse

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement  # 正確匯入

from pptx.enum.text import MSO_AUTO_SIZE
import math

from ..models import ChatRequest

router = APIRouter(prefix="/export", tags=["export"])

# -------------------------
# 字型設定（PDF 用）
# -------------------------
PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent  # ai_agent_backend
FONT_DIR = PROJECT_ROOT / "fonts"
FONT_DIR.mkdir(parents=True, exist_ok=True)

REGULAR_FONT_PATH = FONT_DIR / "NotoSansTC-Regular.ttf"
BOLD_FONT_PATH = FONT_DIR / "NotoSansTC-Bold.ttf"

FONT_NAME = "NotoSansTC"          # 內文字型
FONT_BOLD = "NotoSansTC-Bold"     # 粗體字型


def register_fonts() -> None:
    """
    如果 fonts 資料夾裡有放 NotoSansTC，就註冊進 ReportLab。
    沒有的話，後面會退回用內建 Helvetica（可能會變黑方塊）。
    """
    if REGULAR_FONT_PATH.exists():
        pdfmetrics.registerFont(TTFont(FONT_NAME, str(REGULAR_FONT_PATH)))
    if BOLD_FONT_PATH.exists():
        pdfmetrics.registerFont(TTFont(FONT_BOLD, str(BOLD_FONT_PATH)))


# -------------------------
# 文字前處理：去掉 ### / **...**
# -------------------------
def normalize_line_text(text: str) -> str:
    """
    共用的行文字清理：
    - '### xxx' 變成 ● xxx
    - 移除 Markdown 的 **粗體符號**
    """
    if not text:
        return ""

    # 去掉 Markdown 粗體符號
    t = text.replace("**", "")

    # 處理 '### 標題'
    stripped = t.lstrip()
    if stripped.startswith("### "):
        title = stripped[4:].strip()
        return f"● {title}"

    return t


# -------------------------
# 把對話整理成文字行（PDF 用）
# -------------------------
def flatten_messages(req: ChatRequest) -> List[str]:
    """
    把 ChatRequest 轉成一行一行的文字。
    不在這裡畫 Session ID，Session ID 交給 export_pdf 自己處理。
    """
    lines: List[str] = []

    for m in req.messages:
        if not m.content:
            continue

        if m.role == "user":
            prefix = "我"
        elif m.role == "assistant":
            prefix = "Dr.Bone"
        else:
            prefix = "系統"

        # 保留原本的換行，但先做文字清理
        chunks = m.content.splitlines() or [""]
        for idx, raw_seg in enumerate(chunks):
            seg = normalize_line_text(raw_seg)
            if idx == 0:
                # 第一行帶前綴
                lines.append(f"{prefix}：{seg}")
            else:
                # 後續行縮排但不再重複前綴
                lines.append(f"　{seg}")
        # 每則訊息後補一個空行
        lines.append("")

    return lines


# -------------------------
# PDF：中英文自動換行
# -------------------------
def wrap_text_cjk(
    text: str,
    font_name: str,
    font_size: int,
    max_width: float,
) -> List[str]:
    """
    針對中英混合做簡單換行：
    - 有空白就用空白切（英文句子）
    - 沒空白視為 CJK，一字一字疊加，超出就換行
    """
    if not text:
        return [""]

    # 先觸發一次量測，確保字型已註冊
    pdfmetrics.stringWidth("測", font_name, font_size)

    # 有空白 → 多半是英文 / 中英混合
    if " " in text:
        tokens = text.split(" ")
        space_w = pdfmetrics.stringWidth(" ", font_name, font_size)
        lines: List[str] = []
        cur = ""
        cur_w = 0.0

        for tok in tokens:
            w = pdfmetrics.stringWidth(tok, font_name, font_size)
            if cur and cur_w + space_w + w > max_width:
                lines.append(cur)
                cur = tok
                cur_w = w
            else:
                if cur:
                    cur += " " + tok
                    cur_w += space_w + w
                else:
                    cur = tok
                    cur_w = w

        if cur:
            lines.append(cur)
        return lines

    # 沒有空白 → 當作 CJK，一字一字加
    lines: List[str] = []
    cur = ""
    cur_w = 0.0
    for ch in text:
        w = pdfmetrics.stringWidth(ch, font_name, font_size)
        if cur and cur_w + w > max_width:
            lines.append(cur)
            cur = ch
            cur_w = w
        else:
            cur += ch
            cur_w += w

    if cur:
        lines.append(cur)
    return lines


# -------------------------
# PDF 匯出
# -------------------------
@router.post("/pdf")
def export_pdf(req: ChatRequest):
    """
    匯出 PDF：「骨科互動助理－學習報告」樣式，
    - 標題大一點粗黑字
    - Session ID、一般內文較小
    - 行距依字體大小調整，不會超出邊界
    """
    register_fonts()

    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4

    # 版面設定
    left_margin = 60
    right_margin = 60
    top_margin = 60
    bottom_margin = 60
    usable_width = width - left_margin - right_margin

    # 字型決定（有 Noto 則用，否則退回 Helvetica）
    title_font_name = FONT_BOLD if BOLD_FONT_PATH.exists() else "Helvetica-Bold"
    heading_font_name = FONT_BOLD if BOLD_FONT_PATH.exists() else "Helvetica-Bold"
    body_font_name = FONT_NAME if REGULAR_FONT_PATH.exists() else "Helvetica"

    # 尺寸設定
    title_font_size = 20      # 報告標題
    heading_font_size = 14    # 小標
    body_font_size = 12       # 一般文字

    # ---- 標題 ----
    c.setFont(title_font_name, title_font_size)
    title_text = "骨科互動助理－學習報告"
    c.drawString(left_margin, height - top_margin, title_text)

    # Session ID：前端已經傳「User ID」，直接顯示
    c.setFont(body_font_name, body_font_size)
    y = height - top_margin - (title_font_size + 12)
    c.drawString(left_margin, y, f"Session ID：{req.session_id}")

    # 內文起始位置
    y -= (body_font_size + 12)

    lines = flatten_messages(req)

    for raw_line in lines:
        line = raw_line.rstrip("\n")

        # 空行：直接往下跳一行
        if not line.strip():
            line_spacing = body_font_size + 6
            y -= line_spacing
            if y < bottom_margin:
                c.showPage()
                c.setFont(body_font_name, body_font_size)
                y = height - top_margin
            continue

        # 判斷是不是「一、」「二、」這種小標（如果之後需要）
        stripped = line.lstrip()
        is_heading = stripped.startswith(
            ("一、", "二、", "三、", "四、", "五、", "六、")
        )

        if is_heading:
            cur_font_name = heading_font_name
            cur_font_size = heading_font_size
        else:
            cur_font_name = body_font_name
            cur_font_size = body_font_size

        # 根據目前字型做換行
        wrapped_lines = wrap_text_cjk(
            line, cur_font_name, cur_font_size, usable_width
        )

        for seg in wrapped_lines:
            if y < bottom_margin:
                c.showPage()
                c.setFont(body_font_name, body_font_size)
                y = height - top_margin

            c.setFont(cur_font_name, cur_font_size)
            c.drawString(left_margin, y, seg)
            line_spacing = cur_font_size + 6
            y -= line_spacing

    c.showPage()
    c.save()
    buffer.seek(0)

    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": 'attachment; filename="bone_report.pdf"'},
    )


# -------------------------
# 小工具：把段落的項目符號拿掉（PPTX 用）
# -------------------------
def disable_bullet(paragraph) -> None:
    """
    把這個 paragraph 的 ● 項目符號拿掉。
    """
    pPr = paragraph._element.get_or_add_pPr()
    # 先移除原本的 buChar / buAutoNum / buNone
    for child in list(pPr):
        if child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum") or child.tag.endswith("}buNone"):
            pPr.remove(child)
    bu_none = OxmlElement("a:buNone")
    pPr.insert(0, bu_none)

def remove_all_slides(prs: Presentation) -> None:
    """刪掉簡報中所有既有投影片（保留模板的 layouts/theme）。"""
    sldIdLst = prs.slides._sldIdLst  # type: ignore
    for sldId in list(sldIdLst):
        rId = sldId.rId
        sldIdLst.remove(sldId)
        prs.part.drop_rel(rId)

def wrap_cjk_for_ppt(text: str, max_chars_per_line: int = 32) -> list[str]:
    """超簡單中英混排換行：大概每行最多 N 個字元（WPS/PowerPoint 都吃得到）。"""
    if not text:
        return [""]
    lines = []
    cur = ""
    for ch in text:
        cur += ch
        if len(cur) >= max_chars_per_line:
            lines.append(cur)
            cur = ""
    if cur:
        lines.append(cur)
    return lines

# -------------------------
# PPTX 匯出：標題頁 + 內容頁
# -------------------------
@router.post("/pptx")
def export_pptx(req: ChatRequest):
    TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "templates" / "bone_report_template.pptx"
    print("[export_pptx] template exists =", TEMPLATE_PATH.exists(), "path =", TEMPLATE_PATH)

    prs = Presentation(str(TEMPLATE_PATH)) if TEMPLATE_PATH.exists() else Presentation()

    # ✅ 清掉模板原本的投影片，避免 P1~P3 空白
    remove_all_slides(prs)

    print("[export_pptx] layouts =", len(prs.slide_layouts))

    # ---- 第一頁：標題頁 ----
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = "骨科互動助理－學習報告"
    subtitle_shape.text = f"Session ID：{req.session_id}"

    # 標題字體
    if title_shape.text_frame.paragraphs:
        p = title_shape.text_frame.paragraphs[0]
        for run in p.runs:
            run.font.size = Pt(36)
            run.font.bold = True

    if subtitle_shape.text_frame.paragraphs:
        p = subtitle_shape.text_frame.paragraphs[0]
        for run in p.runs:
            run.font.size = Pt(20)
            run.font.bold = False

    # ---- 第二頁開始：對話內容 ----
    content_layout = prs.slide_layouts[1]  # 標題 + 內容

    def setup_tf(tframe):
        tframe.clear()
        tframe.word_wrap = True
        tframe.margin_left = Inches(0.2)
        tframe.margin_right = Inches(0.2)
        tframe.margin_top = Inches(0.1)
        tframe.margin_bottom = Inches(0.1)
        tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "本次對話內容"
    body = slide.placeholders[1]
    tf = body.text_frame
    setup_tf(tf)

    first_para = True
    max_lines_per_slide = 5
    current_lines = 0

    def add_new_content_slide():
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = "本次對話內容（續）"
        b = s.placeholders[1]
        t = b.text_frame
        setup_tf(t)
        return s, t

    def ensure_space(need_lines: int):
        nonlocal slide, tf, first_para, current_lines
        if current_lines + need_lines > max_lines_per_slide:
            slide, tf = add_new_content_slide()
            first_para = True
            current_lines = 0

    for m in req.messages:
        if not m.content:
            continue

        if m.role == "user":
            prefix = "我"
        elif m.role == "assistant":
            prefix = "Dr.Bone"
        else:
            prefix = "系統"

        chunks = m.content.splitlines() or [""]

        for idx, raw_chunk in enumerate(chunks):
            chunk = normalize_line_text(raw_chunk)

            # ✅ 把內容先用「固定字元數」切成多行，避免爆框
            wrapped = wrap_cjk_for_ppt(chunk, max_chars_per_line=32)
            need_lines = len(wrapped)
            ensure_space(need_lines)

            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()

            p.text = ""
            p.level = 0
            disable_bullet(p)

            # 第一行帶前綴（粗體）
            if idx == 0:
                run_prefix = p.add_run()
                run_prefix.text = f"{prefix}："
                run_prefix.font.size = Pt(20)
                run_prefix.font.bold = (m.role in ("user", "assistant"))

                run_body = p.add_run()
                # 第一行 + 後續行用 \n 串起來（WPS/PowerPoint 都吃）
                run_body.text = wrapped[0] + (
                    "\n" + "\n".join("　" + ln for ln in wrapped[1:]) if len(wrapped) > 1 else ""
                )
                run_body.font.size = Pt(20)
                run_body.font.bold = False

            # 後續行：不帶前綴，純縮排
            else:
                run_body = p.add_run()
                run_body.text = "　" + wrapped[0] + (
                    "\n" + "\n".join("　" + ln for ln in wrapped[1:]) if len(wrapped) > 1 else ""
                )
                run_body.font.size = Pt(20)
                run_body.font.bold = False

            current_lines += need_lines

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)

    return StreamingResponse(
        bio,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="bone_report.pptx"'},
    )
