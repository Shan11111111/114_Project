import os
from pathlib import Path

import pdfplumber
import docx
from pptx import Presentation
import pandas as pd

try:
    from bs4 import BeautifulSoup  # optional
except Exception:
    BeautifulSoup = None


def extract_text_from_pdf(path: str) -> str:
    text_parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text() or ""
            if t:
                text_parts.append(t)
    return "\n".join(text_parts).strip()


def extract_text_from_docx(path: str) -> str:
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs if (p.text or "").strip()]).strip()


def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                s = shape.text.strip()
                if s:
                    parts.append(s)
    return "\n".join(parts).strip()


def extract_text_from_xlsx(path: str) -> str:
    df = pd.read_excel(path)
    return df.to_string(index=False)


def extract_text_from_csv(path: str) -> str:
    df = pd.read_csv(path)
    head = df.head(50)  # 避免爆字
    return head.to_string(index=False)


def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def extract_text_from_md(path: str) -> str:
    return extract_text_from_txt(path)


def extract_text_from_html(path: str) -> str:
    raw = extract_text_from_txt(path)
    if not raw:
        return ""
    if BeautifulSoup:
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text("\n")
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        return "\n".join(lines)
    import re
    text = re.sub(r"<[^>]+>", "\n", raw)
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return "\n".join(lines)


def extract_file_text(path: str) -> str:
    """
    單純做「檔案 → 文字」抽取（不做 RAG / 不做摘要）。
    """
    ext = Path(path).suffix.lower().lstrip(".")

    if ext == "pdf":
        return extract_text_from_pdf(path)
    if ext == "docx":
        return extract_text_from_docx(path)
    if ext == "pptx":
        return extract_text_from_pptx(path)
    if ext in ["xlsx", "xls"]:
        return extract_text_from_xlsx(path)
    if ext == "csv":
        return extract_text_from_csv(path)
    if ext == "txt":
        return extract_text_from_txt(path)
    if ext == "md":
        return extract_text_from_md(path)
    if ext in ["html", "htm"]:
        return extract_text_from_html(path)
    if ext in ["png", "jpg", "jpeg", "webp", "bmp"]:
        return "（圖片檔）建議搭配 YOLO 偵測結果 / BoneDB 映射產生可檢索描述。"

    return "（此檔案格式目前不支援文字解析）"
