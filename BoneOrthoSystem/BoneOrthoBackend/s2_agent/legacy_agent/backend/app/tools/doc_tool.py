# doc_tool.py
from __future__ import annotations

from pathlib import Path
from typing import Tuple, List, Dict, Any, Optional
import re
import os
import json
import hashlib
from urllib.parse import urlparse
from datetime import datetime


# =========================================================
# 你原本的 doc_tool.py（✅ 保留核心行為：clean / summary / extract_* / extract_text_and_summary）
# 另外合併：doc_index（本地 jsonl 索引/檢索）+ URL 抓取與索引
# =========================================================

# -------------------------
# Text helpers
# -------------------------
def _clean_text(t: str) -> str:
    t = (t or "").replace("\r\n", "\n").replace("\r", "\n")
    lines = []
    for line in t.split("\n"):
        line = re.sub(r"[ \t]+", " ", line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines).strip()


def _is_metadata_line(line: str) -> bool:
    """
    過濾上傳檔/網頁常見的行政資訊。
    目標：摘要只留文章重點，不把發布單位、點閱、日期、作者來源塞進去。
    中英文都處理。
    """
    s = re.sub(r"^[\-•●○▪▫◆◇\s]+", "", (line or "").strip())
    if not s:
        return True

    s_lower = s.lower()

    # Excel / 表格標記或純表頭，不應出現在摘要裡
    if re.match(r"^\[?sheet\]?\s*[:：]", s_lower) or re.match(r"^\[sheet\]", s_lower):
        return True
    compact_header = re.sub(r"\s+", " ", s_lower).strip()
    table_headers = [
        "category topic key point details suggested question",
        "test item expected result",
        "life stage goal methods",
    ]
    if compact_header in table_headers:
        return True

    # 日期 / 數字型資訊
    date_patterns = [
        r"\b\d{4}[/-]\d{1,2}[/-]\d{1,2}\b",
        r"\bupdated\s*[:：]",
        r"\bpublished\s*[:：]",
        r"\bviews?\s*[:：]",
        r"更新日期\s*[:：]",
        r"發布日期\s*[:：]",
        r"點閱次數\s*[:：]",
        r"瀏覽次數\s*[:：]",
    ]
    if any(re.search(p, s_lower) for p in date_patterns):
        return True

    # 發布/作者/網站行政資訊
    metadata_phrases = [
        "publishing unit", "published by", "updated", "views",
        "webpage content", "provided by", "listed in order of surname",
        "director of", "department of", "hospital", "medical college",
        "發布單位", "網頁內容", "共同提供", "依姓氏排列",
        "資料來源", "製作單位", "核可", "編號", "copyright", "版權",
    ]
    if any(p in s_lower for p in metadata_phrases):
        return True

    # 純標題且太短，通常不是摘要重點；避免摘要第一條只有「Menopause and Osteoporosis」
    title_like = [
        "menopause and osteoporosis",
        "更年期與骨質疏鬆",
        "risk factors for osteoporosis",
        "how to prevent osteoporosis",
        "骨質疏鬆的危險因子",
        "如何預防骨質疏鬆",
    ]
    if s_lower in title_like:
        return True

    return False


def _drop_boilerplate_lines(lines: List[str]) -> List[str]:
    """
    針對網頁/文件抽字後常見雜訊做基本過濾。
    保留真正衛教內容，移除日期、點閱、發布單位、作者欄位等行政資訊。
    """
    bad_phrases = [
        "熱門文章", "熱門話題", "分類", "標籤", "延伸閱讀",
        "按讚", "分享", "回到首頁", "上一頁", "下一頁",
        "訂閱", "登入", "註冊", "會員", "廣告", "cookie",
        "menu", "home", "login", "subscribe", "share", "advertisement",
    ]
    out = []
    for ln in lines:
        s = ln.strip()
        if not s:
            continue
        if _is_metadata_line(s):
            continue
        # 太短、像導航
        if len(s) <= 2:
            continue
        if any(p.lower() in s.lower() for p in bad_phrases):
            # 但如果是「注意/警訊/就醫」這種關鍵字就留
            if any(k in s for k in ["注意", "警訊", "就醫", "回診", "禁忌", "不可", "避免"]):
                out.append(s)
            continue
        out.append(s)
    return out


def _format_summary_line(line: str) -> str:
    """
    把 Excel 表格列轉成比較像「文章重點」的句子。
    例如：Category: ... | Topic: Spine | Key Point: ... | Details: ...
    會變成：Spine：Osteoporosis commonly affects... Compression fractures...
    """
    s = (line or "").strip()
    if "|" not in s or ":" not in s:
        return s

    fields: Dict[str, str] = {}
    for part in s.split("|"):
        if ":" not in part:
            continue
        k, v = part.split(":", 1)
        k = k.strip().lower()
        v = v.strip()
        if k and v:
            fields[k] = v

    topic = fields.get("topic") or fields.get("主題") or fields.get("項目") or ""
    key_point = fields.get("key point") or fields.get("重點") or fields.get("摘要") or ""
    details = fields.get("details") or fields.get("detail") or fields.get("說明") or fields.get("內容") or ""

    body_parts = []
    if key_point:
        body_parts.append(key_point)
    if details and details not in key_point:
        body_parts.append(details)

    body = " ".join(body_parts).strip()
    if topic and body:
        return f"{topic}: {body}"
    return body or s


def _make_summary(t: str, max_chars: int = 2200) -> str:
    """
    產出文章重點摘要：
    - 中英文文件都用同一套規則
    - 不輸出日期、點閱數、發布單位、作者來源等行政資訊
    - 優先抓疾病/風險/症狀/部位/預防/治療相關句
    """
    t = _clean_text(t)
    if not t:
        return ""

    raw_lines = [ln.strip() for ln in t.split("\n") if ln.strip()]
    if not raw_lines:
        return ""

    lines = _drop_boilerplate_lines(raw_lines)
    if not lines:
        return ""

    keywords = [
        # 中文衛教重點
        "注意", "警訊", "何時", "就醫", "回診", "傷口", "換藥", "浸泡",
        "消毒", "感染", "紅", "腫", "熱", "痛", "流膿", "發燒",
        "步驟", "方法", "需", "請", "避免", "不可", "禁忌",
        "目的", "適應", "適用", "處置", "藥", "敷料",
        "風險", "症狀", "治療", "預防", "檢查", "危險因子",
        "骨質疏鬆", "骨鬆", "更年期", "停經", "雌激素", "骨質流失",
        "骨密度", "骨折", "脊椎", "壓迫性骨折", "股骨", "髖骨", "手腕",
        "鈣", "維他命", "維生素", "運動", "抽菸", "吸菸", "酗酒", "類固醇",

        # English key points
        "key point", "details", "risk", "risk factor", "symptom", "treatment", "therapy", "prevention", "prevent",
        "screening", "check", "diagnosis", "osteoporosis", "osteopenia", "menopause",
        "postmenopausal", "estrogen", "bone loss", "bone density", "fracture",
        "compression fracture", "spine", "vertebra", "vertebrae", "hip", "femoral",
        "femur", "wrist", "calcium", "vitamin d", "exercise", "smoking", "alcohol",
        "corticosteroid", "steroid", "diabetes", "kidney disease", "dialysis",
    ]

    picked: List[str] = []
    seen = set()

    def normalize_for_seen(line: str) -> str:
        return re.sub(r"\s+", " ", line.strip().lower())

    def add(line: str):
        s = _format_summary_line(line.strip())
        if not s or _is_metadata_line(s):
            return
        # 避免把表格殘片或單一欄位塞進摘要
        if len(s) < 12 and not any(k in s for k in ["骨折", "骨鬆", "骨質疏鬆"]):
            return
        key = normalize_for_seen(s)
        if key in seen:
            return
        seen.add(key)
        picked.append(s)

    # 1) 優先抓含中英文關鍵詞的句子
    for ln in lines:
        low = ln.lower()
        if any(k.lower() in low for k in keywords):
            add(ln)
        if len(picked) >= 14:
            break

    # 2) 不足的話，補前段有資訊量的內容，但仍過濾行政資訊
    if len(picked) < 6:
        for ln in lines[:40]:
            add(ln)
            if len(picked) >= 10:
                break

    out = []
    total = 0
    for ln in picked:
        item = f"- {ln}"
        if total + len(item) > max_chars:
            break
        out.append(item)
        total += len(item)

    return "\n".join(out).strip()


# -------------------------
# Extractors
# -------------------------
def _extract_pdf(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(path))
        texts = []
        for page in doc:
            texts.append(page.get_text("text") or "")
        doc.close()
        return "\n".join(texts)
    except Exception:
        pass

    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(str(path)) as pdf:
            for p in pdf.pages:
                texts.append(p.extract_text() or "")
        return "\n".join(texts)
    except Exception:
        pass

    try:
        from pypdf import PdfReader
        r = PdfReader(str(path))
        texts = []
        for p in r.pages:
            texts.append(p.extract_text() or "")
        return "\n".join(texts)
    except Exception as e:
        raise RuntimeError(f"PDF 解析器不可用或解析失敗：{e}")


def _extract_txt_like(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _extract_docx(path: Path) -> str:
    try:
        import docx  # python-docx
        d = docx.Document(str(path))
        return "\n".join([p.text for p in d.paragraphs])
    except Exception as e:
        raise RuntimeError(f"DOCX 解析失敗：{e}")


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text or "")
        return "\n".join(texts)
    except Exception as e:
        raise RuntimeError(f"PPTX 解析失敗：{e}")


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True)
        texts = []

        for ws in wb.worksheets:
            texts.append(f"[Sheet] {ws.title}")

            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            # 找第一列有內容的 row 當表頭；如果不像表頭，仍用一般文字方式保底
            header_row = None
            header_idx = -1
            for idx, row in enumerate(rows):
                vals = [str(x).strip() for x in row if x is not None and str(x).strip()]
                if vals:
                    header_row = [str(x).strip() if x is not None else "" for x in row]
                    header_idx = idx
                    break

            if header_row is None:
                continue

            headers = [h.strip() for h in header_row]
            has_headers = sum(1 for h in headers if h) >= 2

            for row in rows[header_idx + 1 if has_headers else header_idx:]:
                cells = [str(x).strip() if x is not None else "" for x in row]
                if not any(cells):
                    continue

                if has_headers:
                    parts = []
                    for h, c in zip(headers, cells):
                        if h and c:
                            parts.append(f"{h}: {c}")
                    # 如果欄位數超過表頭，也把剩下有內容的 cell 留下
                    if len(cells) > len(headers):
                        parts.extend([c for c in cells[len(headers):] if c])
                    line = " | ".join(parts).strip()
                else:
                    line = " ".join([c for c in cells if c]).strip()

                if line:
                    texts.append(line)

        return "\n".join(texts)
    except Exception as e:
        raise RuntimeError(f"XLSX 解析失敗：{e}")


def _strip_html(html: str) -> str:
    s = html or ""
    s = re.sub(r"(?is)<(script|style|noscript|svg|canvas|iframe).*?>.*?</\1>", " ", s)
    s = re.sub(r"(?is)<!--.*?-->", " ", s)
    s = re.sub(r"(?is)<br\s*/?>", "\n", s)
    s = re.sub(r"(?is)</p\s*>", "\n", s)
    s = re.sub(r"(?is)</div\s*>", "\n", s)
    s = re.sub(r"(?is)<[^>]+>", " ", s)
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def _extract_html_like(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    return _strip_html(raw)


def _extract_md(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    raw = re.sub(r"```.*?```", " ", raw, flags=re.S)
    raw = re.sub(r"`([^`]+)`", r"\1", raw)
    raw = re.sub(r"!\[[^\]]*\]\([^)]+\)", " ", raw)
    raw = re.sub(r"\[[^\]]+\]\([^)]+\)", " ", raw)
    raw = re.sub(r"^#{1,6}\s*", "", raw, flags=re.M)
    return raw


# -------------------------
# Public API (原本的)
# -------------------------
def extract_text_and_summary(file_path: Path, ext: str) -> Tuple[str, str]:
    p = Path(file_path)
    e = (ext or "").lower().lstrip(".")

    if not p.exists():
        raise RuntimeError(f"檔案不存在：{p}")

    if e == "pdf":
        text = _extract_pdf(p)
    elif e in ("txt", "csv"):
        text = _extract_txt_like(p)
    elif e in ("doc", "docx"):
        if e == "doc":
            raise RuntimeError("不支援 .doc（請轉成 .docx）")
        text = _extract_docx(p)
    elif e in ("ppt", "pptx"):
        if e == "ppt":
            raise RuntimeError("不支援 .ppt（請轉成 .pptx）")
        text = _extract_pptx(p)
    elif e in ("xls", "xlsx"):
        if e == "xls":
            raise RuntimeError("不支援 .xls（請轉成 .xlsx）")
        text = _extract_xlsx(p)
    elif e in ("html", "htm"):
        text = _extract_html_like(p)
    elif e in ("md", "markdown"):
        text = _extract_md(p)
    else:
        raise RuntimeError(f"不支援的檔案格式：.{e}")

    text = _clean_text(text)
    if not text:
        raise RuntimeError("抽不到文字：可能是掃描 PDF（沒有文字層），或檔案內容不可抽取。")

    summary = _make_summary(text)
    return text, summary


# =========================================================
# ✅ Doc-RAG：本地索引/檢索（jsonl）
#   - 開關：S2_ENABLE_DOC_RAG=1
# =========================================================
TOOLS_DIR = Path(__file__).resolve().parent
INDEX_DIR = TOOLS_DIR / "_doc_index"
INDEX_DIR.mkdir(parents=True, exist_ok=True)
INDEX_PATH = INDEX_DIR / "doc_store.jsonl"


def is_enabled() -> bool:
    return os.getenv("S2_ENABLE_DOC_RAG", "0") == "1"


def _normalize_token(tok: str) -> str:
    """
    很輕量的 token 正規化：
    - 英文轉小寫
    - 常見複數/時態做保守處理，避免 bones/fractures 跟 bone/fracture 完全對不上
    - 中文仍維持單字切分，避免破壞你原本的中文 overlap 行為
    """
    t = (tok or "").strip().lower()
    if not t:
        return ""
    if re.fullmatch(r"[a-z0-9]+", t):
        if len(t) > 5 and t.endswith("ies"):
            t = t[:-3] + "y"
        elif len(t) > 4 and t.endswith("es"):
            t = t[:-2]
        elif len(t) > 3 and t.endswith("s"):
            t = t[:-1]
        elif len(t) > 5 and t.endswith("ed"):
            t = t[:-2]
        elif len(t) > 6 and t.endswith("ing"):
            t = t[:-3]
    return t


def _tokenize(s: str) -> List[str]:
    s = (s or "").lower()
    parts = re.findall(r"[\u4e00-\u9fff]|[a-z0-9]+", s)
    out: List[str] = []
    for p in parts:
        t = _normalize_token(p)
        if t:
            out.append(t)
    return out


def _overlap_score(q_tokens: List[str], d_tokens: List[str]) -> float:
    if not q_tokens or not d_tokens:
        return 0.0
    qs = set(q_tokens)
    ds = set(d_tokens)
    inter = len(qs & ds)
    denom = (len(qs) ** 0.7) * (len(ds) ** 0.3)
    return float(inter) / float(denom or 1.0)



# ---------------------------------------------------------
# ✅ Bilingual query expansion for Doc-RAG
#   目的：
#   - 使用者用中文問，但上傳文件是英文時，仍能撈到內容
#   - 使用者用英文問，但資料是中文時，也補常見中文詞
#   - 不取代向量資料庫，只補強你目前 jsonl overlap 檢索的弱點
# ---------------------------------------------------------
_DOC_QUERY_SYNONYMS: Dict[str, str] = {
    # osteoporosis / menopause
    "骨質疏鬆": "osteoporosis osteopenia low bone density bone loss 骨鬆 骨密度 骨質流失",
    "骨质疏松": "osteoporosis osteopenia low bone density bone loss 骨鬆 骨密度 骨質流失",
    "骨鬆": "osteoporosis osteopenia low bone density bone loss 骨質疏鬆 骨密度",
    "骨密度": "bone mineral density bmd low bone density osteoporosis 骨質疏鬆",
    "骨質流失": "bone loss osteoporosis estrogen menopause postmenopausal 骨質疏鬆",
    "更年期": "menopause menopausal postmenopausal estrogen hormone hormone replacement therapy 停經 雌激素",
    "停經": "menopause menopausal postmenopausal estrogen hormone hormone replacement therapy 更年期 雌激素",
    "停經後": "postmenopausal after menopause estrogen bone loss osteoporosis",
    "雌激素": "estrogen oestrogen female hormone menopause postmenopausal bone loss",
    "荷爾蒙": "hormone estrogen hormone replacement therapy hrt menopause",

    # anatomy / fracture sites
    "脊椎": "spine vertebra vertebrae vertebral spinal compression fracture thoracic lumbar cervical 脊椎骨 椎骨",
    "脊椎骨": "spine vertebra vertebrae vertebral spinal compression fracture",
    "椎骨": "vertebra vertebrae spine vertebral compression fracture",
    "頸椎": "cervical spine cervical vertebrae neck vertebra",
    "胸椎": "thoracic spine thoracic vertebrae hunchback kyphosis",
    "腰椎": "lumbar spine lumbar vertebrae lower back pain",
    "肋骨": "rib ribs rib cage costal bone thoracic cage",
    "骨盆": "pelvis pelvic bone hip",
    "髖骨": "hip pelvis femoral head hip fracture",
    "股骨": "femur femoral head thigh bone hip fracture",
    "大腿骨": "femur femoral head thigh bone hip fracture",
    "手腕": "wrist carpal bones wrist fracture distal radius",
    "腕骨": "wrist carpal bones wrist fracture",
    "骨折": "fracture broken bone compression fracture hip fracture wrist fracture",
    "壓迫性骨折": "compression fracture vertebral compression fracture spine osteoporosis",

    # risk factors / prevention
    "鈣": "calcium dairy milk calcium intake",
    "維他命d": "vitamin d vitamin d deficiency calcium bone density",
    "維生素d": "vitamin d vitamin d deficiency calcium bone density",
    "運動": "exercise physical activity weight bearing exercise",
    "抽菸": "smoking tobacco risk factor osteoporosis",
    "吸菸": "smoking tobacco risk factor osteoporosis",
    "酗酒": "alcohol alcohol abuse risk factor osteoporosis",
    "類固醇": "corticosteroid steroid glucocorticoid medication bone loss osteoporosis",
    "糖尿病": "diabetes endocrine disease osteoporosis risk factor",
    "洗腎": "dialysis kidney disease renal disease osteoporosis",
    "預防": "prevention prevent calcium vitamin d exercise screening",
    "治療": "treatment therapy hormone replacement therapy calcitonin fluoride medication",

    # English -> Chinese補強
    "osteoporosis": "骨質疏鬆 骨鬆 骨密度 bone loss low bone density",
    "osteopenia": "骨質不足 骨密度 low bone density",
    "menopause": "更年期 停經 postmenopausal estrogen",
    "postmenopausal": "停經後 更年期 estrogen menopause",
    "estrogen": "雌激素 荷爾蒙 menopause bone loss",
    "spine": "脊椎 椎骨 vertebrae compression fracture",
    "vertebra": "椎骨 脊椎 vertebrae spine",
    "vertebrae": "椎骨 脊椎 vertebra spine",
    "rib": "肋骨 ribs",
    "ribs": "肋骨 rib",
    "pelvis": "骨盆 pelvic hip",
    "hip": "髖骨 股骨 femoral head hip fracture",
    "femur": "股骨 大腿骨 femoral head",
    "femoral": "股骨 femur femoral head",
    "wrist": "手腕 腕骨 carpal",
    "fracture": "骨折 compression fracture hip fracture wrist fracture",
    "calcium": "鈣 鈣質 dairy milk",
    "vitamin": "維他命 維生素 vitamin d",
    "corticosteroid": "類固醇 steroid glucocorticoid",
    "steroid": "類固醇 corticosteroid glucocorticoid",
    "diabetes": "糖尿病 endocrine disease",
}


def expand_doc_query(query: str) -> str:
    """
    給 doc_tool.retrieve() 使用的查詢擴展。
    這裡故意做成純規則，不碰 OpenAI、不碰 Qdrant，避免你原本本地 doc index 被改爆。
    """
    q = (query or "").strip()
    if not q:
        return ""

    q_lower = q.lower()
    extras: List[str] = []

    for key, val in _DOC_QUERY_SYNONYMS.items():
        k = key.lower()
        if k in q_lower:
            extras.append(val)

    # 使用者只問「前面那個 / 這個 / 它」時，通常沒有足夠關鍵字。
    # 這裡只補泛用骨科詞，避免完全查不到；真正上下文仍應由 main/rag_tool 把歷史問題合併進來。
    vague_words = ["前面", "剛剛", "這個", "那個", "它", "this", "that", "it"]
    if any(w in q_lower for w in vague_words):
        extras.append("bone anatomy osteoporosis fracture spine rib hip wrist menopause")

    if not extras:
        return q

    # 去重但保留順序
    seen = set()
    merged: List[str] = []
    for item in extras:
        for tok in item.split():
            t = tok.strip()
            if t and t.lower() not in seen:
                seen.add(t.lower())
                merged.append(t)

    return (q + "\n" + " ".join(merged)).strip()


def _split_paragraphs(text: str) -> List[str]:
    t = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    paras = re.split(r"\n\s*\n", t)
    out = []
    for p in paras:
        p = _clean_text(p)
        if p:
            out.append(p)
    return out


def _chunk_text(text: str, max_chars: int = 900) -> List[str]:
    paras = _split_paragraphs(text)
    chunks: List[str] = []
    buf = ""
    for p in paras:
        if not buf:
            buf = p
            continue
        if len(buf) + 1 + len(p) <= max_chars:
            buf = buf + "\n" + p
        else:
            chunks.append(buf)
            buf = p
    if buf:
        chunks.append(buf)
    return chunks


def _append_jsonl(path: Path, rows: List[Dict[str, Any]]) -> None:
    if not rows:
        return
    with open(path, "a", encoding="utf-8") as f:
        for r in rows:
            f.write(json.dumps(r, ensure_ascii=False) + "\n")


def _read_all_jsonl(path: Path) -> List[Dict[str, Any]]:
    if not path.exists():
        return []
    out = []
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                out.append(json.loads(line))
            except Exception:
                continue
    return out


def index_document(
    text: str,
    title: str,
    source_type: str,
    material_id: str,
    *,
    url: Optional[str] = None,
    conversation_id: Optional[str] = None,
    user_id: Optional[str] = None,
) -> int:
    if not text:
        return 0
    chunks = _chunk_text(text, max_chars=int(os.getenv("S2_DOC_CHUNK_CHARS", "900")))
    rows = []
    now = datetime.utcnow().isoformat()

    for i, ch in enumerate(chunks):
        rows.append(
            {
                "material_id": str(material_id),
                "title": title,
                "source_type": source_type,
                "chunk_index": i,
                "text": ch,
                # 把 title 一起放進 tokens，讓「問檔名/主題」時也比較容易命中。
                "tokens": _tokenize(ch + "\n" + (title or "")),
                "url": url,
                "conversation_id": conversation_id,
                "user_id": user_id,
                "created_at": now,
            }
        )

    _append_jsonl(INDEX_PATH, rows)
    return len(rows)


def retrieve(query: str, top_k: int = 6) -> List[Dict[str, Any]]:
    q = (query or "").strip()
    if not q:
        return []

    # ✅ 重要：先做中英雙語擴展
    # 例：使用者問「肋骨與骨質疏鬆」，英文文件裡是 ribs / osteoporosis，原本 overlap 會撈不到。
    expanded_q = expand_doc_query(q)
    q_tokens = _tokenize(expanded_q)
    docs = _read_all_jsonl(INDEX_PATH)

    scored: List[Dict[str, Any]] = []
    for d in docs:
        d_tokens = d.get("tokens") or []
        s = _overlap_score(q_tokens, d_tokens)

        # title 命中時給一點點加權，避免同主題文件分數太低。
        title = d.get("title") or ""
        title_score = _overlap_score(q_tokens, _tokenize(title))
        if title_score > 0:
            s += min(0.12, title_score * 0.25)

        if s <= 0:
            continue

        scored.append(
            {
                "material_id": d.get("material_id"),
                "title": d.get("title"),
                "source_type": d.get("source_type"),
                "chunk_index": d.get("chunk_index"),
                "text": d.get("text"),
                "score": float(s),
                "url": d.get("url"),
            }
        )

    scored.sort(key=lambda x: x["score"], reverse=True)

    # 避免同一個 chunk 因為重複索引被一直回傳。
    deduped: List[Dict[str, Any]] = []
    seen_keys = set()
    for item in scored:
        key = (
            item.get("material_id"),
            item.get("chunk_index"),
            (item.get("text") or "")[:80],
        )
        if key in seen_keys:
            continue
        seen_keys.add(key)
        deduped.append(item)
        if len(deduped) >= max(1, int(top_k)):
            break

    return deduped


# =========================================================
# ✅ URL 抓取 + 索引（貼 URL 就能解析）
#   - 403 / 斷線：不要讓外層爆炸（index_url 不丟例外）
#   - 加 headers + retry + jina.ai fallback
# =========================================================
_URL_RE = re.compile(r"https?://\S+")


def extract_first_url(text: str) -> Optional[str]:
    m = _URL_RE.search(text or "")
    if not m:
        return None
    return m.group(0).strip().rstrip(").,]}>\"'")


def _guess_title_from_html(html: str) -> str:
    m = re.search(r"(?is)<title[^>]*>(.*?)</title>", html or "")
    if not m:
        return ""
    t = re.sub(r"\s+", " ", m.group(1)).strip()
    return t[:120]


def _download_url(url: str) -> Tuple[str, str, str]:
    """
    return: (body_text, final_url, content_type)
    """
    try:
        import httpx
    except Exception as e:
        raise RuntimeError(f"httpx 未安裝：{e}")

    headers = {
        "User-Agent": os.getenv(
            "S2_URL_UA",
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36",
        ),
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "zh-TW,zh;q=0.9,en;q=0.7",
        "Cache-Control": "no-cache",
        "Pragma": "no-cache",
        "Connection": "close",
    }

    timeout = float(os.getenv("S2_URL_TIMEOUT", "12"))
    max_bytes = int(os.getenv("S2_URL_MAX_BYTES", "2500000"))  # 2.5MB

    last_err = None
    for _ in range(2):  # retry
        try:
            with httpx.Client(
                follow_redirects=True,
                timeout=timeout,
                headers=headers,
                http2=False,
            ) as client:
                r = client.get(url)
                ct = (r.headers.get("content-type") or "").lower()
                if r.status_code >= 400:
                    raise RuntimeError(f"Client error '{r.status_code}' for url '{str(r.url)}'")
                text = r.text
                if len(text.encode("utf-8", errors="ignore")) > max_bytes:
                    text = text[: max_bytes // 2]  # 粗暴截斷，避免爆記憶體/延遲
                return text, str(r.url), ct
        except Exception as e:
            last_err = e
            continue

    raise RuntimeError(str(last_err))


def _download_url_via_jina(url: str) -> Tuple[str, str, str]:
    """
    fallback：把 https://xxx 變成 https://r.jina.ai/https://xxx
    403 有時會被救起來（非保證；付費牆/強反爬仍可能失敗）
    """
    try:
        import httpx
    except Exception as e:
        raise RuntimeError(f"httpx 未安裝：{e}")

    u = url.strip()
    if u.startswith("https://"):
        proxy = "https://r.jina.ai/https://" + u[len("https://") :]
    elif u.startswith("http://"):
        proxy = "https://r.jina.ai/http://" + u[len("http://") :]
    else:
        proxy = "https://r.jina.ai/http://" + u

    timeout = float(os.getenv("S2_URL_TIMEOUT", "12"))
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept": "text/plain,*/*",
        "Connection": "close",
    }
    with httpx.Client(follow_redirects=True, timeout=timeout, headers=headers, http2=False) as client:
        r = client.get(proxy)
        if r.status_code >= 400:
            raise RuntimeError(f"Proxy error '{r.status_code}' for url '{proxy}'")
        text = r.text
        return text, url, "text/plain"


def extract_url_text_and_summary(url: str) -> Tuple[str, str, str, str]:
    """
    return: (text, summary, title, final_url)
    """
    body = ""
    final_url = url
    ct = ""
    last_err = None

    try:
        body, final_url, ct = _download_url(url)
    except Exception as e:
        last_err = e

    if not body:
        try:
            body, final_url, ct = _download_url_via_jina(url)
            last_err = None
        except Exception as e:
            last_err = e

    if not body:
        raise RuntimeError(f"下載網址失敗：{last_err}")

    title = ""
    if "text/html" in ct or "<html" in body.lower():
        title = _guess_title_from_html(body)
        text = _strip_html(body)
    else:
        text = body

    text = _clean_text(text)
    if not text:
        raise RuntimeError("抽不到文字：該網址可能是純圖片/需登入/或拒絕存取。")

    # 再做一輪雜訊過濾（避免摘要全是導覽）
    lines = _drop_boilerplate_lines([ln for ln in text.split("\n") if ln.strip()])
    text = "\n".join(lines).strip()

    summary = _make_summary(text)
    if not title:
        try:
            host = urlparse(final_url).netloc
            title = host or "web"
        except Exception:
            title = "web"

    return text, summary, title, final_url


def index_url(
    url: str,
    conversation_id: Optional[str] = None,
    user_id: Optional[str] = None,
) -> Dict[str, Any]:
    """
    ✅ 最重要：這個函式「不要把錯丟出去」
    成功/失敗都回 dict，讓 main.py 能把結果塞進 (1) 直接回答前面的 context
    """
    url = (url or "").strip()
    if not url:
        return {"ok": False, "url": url, "warning": "empty url"}

    try:
        text, summary, title, final_url = extract_url_text_and_summary(url)

        mid = hashlib.sha1(final_url.encode("utf-8")).hexdigest()
        indexed = index_document(
            text=text,
            title=title,
            source_type="url",
            material_id=mid,
            url=final_url,
            conversation_id=conversation_id,
            user_id=user_id,
        )

        return {
            "ok": True,
            "material_id": mid,
            "title": title,
            "summary": summary,
            "text": text,
            "url": final_url,
            "indexed_chunks": indexed,
        }
    except Exception as e:
        return {
            "ok": False,
            "url": url,
            "warning": f"{e}",
            "material_id": hashlib.sha1(url.encode('utf-8')).hexdigest(),
        }


def build_url_digest(idx: Dict[str, Any]) -> str:
    """
    給 main.py / rag prompt 用：
    讓「已解析網址/摘要」能放到 (1) 直接回答之前
    """
    if not idx:
        return ""

    ok = bool(idx.get("ok"))
    url = idx.get("url") or ""
    title = idx.get("title") or ""
    summary = (idx.get("summary") or "").strip()
    warning = (idx.get("warning") or "").strip()

    if ok:
        head = f"【已解析網址】{url}"
        if title:
            head += f"\n【標題】{title}"
        if summary:
            return f"{head}\n【摘要】\n{summary}".strip()
        return f"{head}\n【摘要】（抽取成功，但摘要為空）".strip()

    # 失敗：也要回一段可以讓模型/使用者理解的文字，不要直接消失
    head = f"【已解析網址】{url}"
    if warning:
        return f"{head}\n【解析失敗】{warning}".strip()
    return f"{head}\n【解析失敗】未知錯誤".strip()
