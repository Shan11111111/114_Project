# BoneOrthoBackend/s2_agent/vectordb/ingest_materials.py
from dotenv import load_dotenv
from pathlib import Path
from typing import Dict, Any, List, Union
import json
import re
import hashlib
from shared.vector_client import VectorStore, is_gibberish

# ----------------------------
# Load .env (BoneOrthoBackend/.env)
# ----------------------------
ENV_CANDIDATES = [
    Path(__file__).resolve().parents[2] / ".env",  # BoneOrthoBackend/.env
    Path(__file__).resolve().parents[1] / ".env",
    Path(__file__).resolve().parents[3] / ".env",
]
for p in ENV_CANDIDATES:
    if p.exists():
        load_dotenv(p)
        break

from db import get_connection
from shared.vector_client import VectorStore
from s2_agent.service import embed_text

# 指向 BoneOrthoBackend/s2_agent/vectordb/materials
MATERIAL_ROOT = Path(__file__).resolve().parent / "materials"

# ----------------------------
# Default StructureJson config
# ----------------------------
DEFAULT_CONFIG: Dict[str, Any] = {
    "version": 1,
    "loader": "auto",
    "splitter": "recursive",
    "chunk_size": 1000,
    "chunk_overlap": 150,
    "min_chunk_len": 30,
    "ocr": False,
    "keep_metadata": True,
    "keep_title": True,
    "keep_page_number": True,
    "keep_slide_number": True,
    "rows_per_chunk": 20,
    "sheet_mode": "all",
    "sheet_name": None,
    "json_path": None,
    "language_hint": "mixed",
}

SUPPORTED_LOADERS = {"auto", "pdf", "txt", "note", "md", "docx", "pptx"}
SUPPORTED_SPLITTERS = {"recursive", "paragraph", "slide"}


# ----------------------------
# Utilities
# ----------------------------
def _normalize_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\x00", " ")
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _safe_material_path(rel_path: str) -> Path:
    """
    只允許 materials 底下的相對路徑，避免路徑穿越。
    """
    rel = Path(str(rel_path or "").replace("\\", "/"))
    if rel.is_absolute():
        raise ValueError(f"FilePath 不可為絕對路徑: {rel_path}")

    full = (MATERIAL_ROOT / rel).resolve()
    root = MATERIAL_ROOT.resolve()

    if root not in full.parents and full != root:
        raise ValueError(f"FilePath 超出 materials 目錄: {rel_path}")

    return full


import uuid
import hashlib

def _make_doc_id(material_id: str, chunk_index: int, text: str) -> str:
    seed = f"{material_id}|{chunk_index}|{hashlib.md5(text.encode('utf-8')).hexdigest()}"
    return str(uuid.uuid5(uuid.NAMESPACE_URL, seed))


def _split_recursive(text: str, config: Dict[str, Any]) -> List[str]:
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    text = _normalize_text(text)
    if not text:
        return []

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=int(config.get("chunk_size", 1000)),
        chunk_overlap=int(config.get("chunk_overlap", 150)),
        separators=["\n\n", "\n", "。", "！", "？", ".", " ", ""],
    )
    parts = splitter.split_text(text)
    min_len = int(config.get("min_chunk_len", 30))
    return [p.strip() for p in parts if p and len(p.strip()) >= min_len]


def _split_paragraph(text: str, config: Dict[str, Any]) -> List[str]:
    text = _normalize_text(text)
    if not text:
        return []

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    min_len = int(config.get("min_chunk_len", 30))
    chunk_size = int(config.get("chunk_size", 1000))
    chunk_overlap = int(config.get("chunk_overlap", 150))

    results: List[str] = []
    buf = ""

    for p in paragraphs:
        if len(p) < min_len:
            continue

        candidate = p if not buf else f"{buf}\n\n{p}"
        if len(candidate) <= chunk_size:
            buf = candidate
        else:
            if buf:
                results.append(buf.strip())
            buf = p

    if buf:
        results.append(buf.strip())

    # 若段落合併後仍太長，再用 recursive 補切
    final_parts: List[str] = []
    for item in results:
        if len(item) > chunk_size:
            final_parts.extend(_split_recursive(item, config))
        else:
            final_parts.append(item)

    # 簡單 overlap 補強：若需要可對 paragraph 再做更精細處理
    if chunk_overlap > 0:
        cleaned = [x for x in final_parts if len(x.strip()) >= min_len]
        return cleaned

    return [x for x in final_parts if len(x.strip()) >= min_len]


def _resolve_config(row: Dict[str, Any]) -> Dict[str, Any]:
    cfg = dict(DEFAULT_CONFIG)
    raw = row.get("StructureJson")

    if raw:
        try:
            user_cfg = raw if isinstance(raw, dict) else json.loads(raw)
            if isinstance(user_cfg, dict):
                cfg.update({k: v for k, v in user_cfg.items() if v is not None})
        except Exception as e:
            print(f"⚠️ StructureJson parse failed, use default config: {e}")

    loader = str(cfg.get("loader", "auto")).lower().strip()
    splitter = str(cfg.get("splitter", "recursive")).lower().strip()

    if loader not in SUPPORTED_LOADERS:
        print(f"⚠️ unsupported loader={loader}, fallback to auto")
        cfg["loader"] = "auto"

    if splitter not in SUPPORTED_SPLITTERS:
        print(f"⚠️ unsupported splitter={splitter}, fallback to recursive")
        cfg["splitter"] = "recursive"

    return cfg


def _resolve_loader(row: Dict[str, Any], full_path: Path, config: Dict[str, Any]) -> str:
    loader = str(config.get("loader", "auto")).lower().strip()
    if loader != "auto":
        return loader

    typ = str(row.get("Type") or "").lower().strip()
    if typ in SUPPORTED_LOADERS:
        return typ

    suffix = full_path.suffix.lower()
    mapping = {
        ".pdf": "pdf",
        ".txt": "txt",
        ".md": "md",
        ".docx": "docx",
        ".pptx": "pptx",
    }
    return mapping.get(suffix, typ or "txt")


def _apply_splitter(text: str, config: Dict[str, Any]) -> List[str]:
    splitter = str(config.get("splitter", "recursive")).lower().strip()

    if splitter == "paragraph":
        return _split_paragraph(text, config)

    # slide 這個 splitter 通常在 pptx loader 內處理；若落到純文字，仍用 recursive
    return _split_recursive(text, config)


# ----------------------------
# DB
# ----------------------------
def _fetch_material(material_id: Union[str, bytes]) -> Dict[str, Any]:
    sql = """
    SELECT MaterialId, Type, Language, Style, Title,
           StructureJson, FilePath, CreatedAt
    FROM agent.TeachingMaterial
    WHERE MaterialId = ?
    """
    with get_connection() as conn:
        cur = conn.cursor()
        cur.execute(sql, material_id)
        row = cur.fetchone()
        if not row:
            raise ValueError(f"TeachingMaterial {material_id} not found")
        cols = [c[0] for c in cur.description]

    data = dict(zip(cols, row))

    # 補預設欄位，避免後面程式 row.get("BoneId") / row.get("BoneSmallId") 炸掉
    data.setdefault("BoneId", None)
    data.setdefault("BoneSmallId", None)

    return data


# ----------------------------
# Loader implementations
# ----------------------------
def _load_pdf_chunks(full_path: Path, row: Dict[str, Any], config: Dict[str, Any]) -> List[Dict[str, Any]]:
    docs = []

    try:
        from langchain_community.document_loaders import PyMuPDFLoader
        loader = PyMuPDFLoader(str(full_path))
        docs = loader.load()
    except Exception as e:
        print(f"⚠️ PyMuPDFLoader failed, fallback PyPDFLoader: {e}")

    if not docs:
        try:
            from langchain_community.document_loaders import PyPDFLoader
            loader = PyPDFLoader(str(full_path))
            docs = loader.load()
        except Exception as e:
            print(f"❌ PyPDFLoader also failed: {e}")
            return []

    chunks: List[Dict[str, Any]] = []
    chunk_index = 0

    for d in docs:
        raw_text = getattr(d, "page_content", None) or ""
        raw_text = _normalize_text(raw_text)
        if not raw_text:
            continue

        meta = getattr(d, "metadata", {}) or {}
        page_num = meta.get("page", None)

        if isinstance(page_num, int) and bool(config.get("keep_page_number", True)):
            page_num = page_num + 1
        else:
            page_num = None

        parts = _apply_splitter(raw_text, config)
        for part in parts:
            chunks.append({
                "text": part,
                "page": page_num if bool(config.get("keep_metadata", True)) else None,
                "slide": None,
                "chunk_index": chunk_index,
            })
            chunk_index += 1

    if not chunks:
        print(f"⚠️ PDF 無可用文字內容，略過索引: {row.get('Title')} ({row.get('FilePath')})")

    return chunks


def _load_text_like_chunks(full_path: Path, row: Dict[str, Any], config: Dict[str, Any]) -> List[Dict[str, Any]]:
    text = full_path.read_text(encoding="utf-8", errors="ignore")
    text = _normalize_text(text)
    if not text:
        print(f"⚠️ empty text file, skip: {row.get('Title')} ({row.get('FilePath')})")
        return []

    parts = _apply_splitter(text, config)
    return [
        {
            "text": part,
            "page": None,
            "slide": None,
            "chunk_index": i,
        }
        for i, part in enumerate(parts)
    ]


def _load_docx_chunks(full_path: Path, row: Dict[str, Any], config: Dict[str, Any]) -> List[Dict[str, Any]]:
    try:
        from docx import Document
    except ImportError:
        print("❌ python-docx 未安裝，無法處理 docx")
        return []

    try:
        doc = Document(str(full_path))
    except Exception as e:
        print(f"❌ docx load failed: {e}")
        return []

    paragraphs = []
    for p in doc.paragraphs:
        txt = _normalize_text(p.text)
        if txt:
            paragraphs.append(txt)

    if not paragraphs:
        print(f"⚠️ docx 無文字內容，略過索引: {row.get('Title')}")
        return []

    whole_text = "\n\n".join(paragraphs)
    parts = _apply_splitter(whole_text, config)

    return [
        {
            "text": part,
            "page": None,
            "slide": None,
            "chunk_index": i,
        }
        for i, part in enumerate(parts)
    ]


def _load_pptx_chunks(full_path: Path, row: Dict[str, Any], config: Dict[str, Any]) -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError:
        print("❌ python-pptx 未安裝，無法處理 pptx")
        return []

    try:
        prs = Presentation(str(full_path))
    except Exception as e:
        print(f"❌ pptx load failed: {e}")
        return []

    keep_title = bool(config.get("keep_title", True))
    keep_metadata = bool(config.get("keep_metadata", True))
    splitter = str(config.get("splitter", "slide")).lower().strip()

    chunks: List[Dict[str, Any]] = []
    chunk_index = 0

    for idx, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = _normalize_text(shape.text)
                if txt:
                    texts.append(txt)

        if not texts:
            continue

        slide_text = "\n".join(texts).strip()
        if not slide_text:
            continue

        title_text = ""
        if keep_title and getattr(slide.shapes, "title", None) is not None:
            try:
                title_text = _normalize_text(slide.shapes.title.text or "")
            except Exception:
                title_text = ""

        if title_text and title_text not in slide_text:
            slide_text = f"{title_text}\n\n{slide_text}"

        # slide splitter：每張投影片當主單位；若太長再 recursive 細切
        if splitter == "slide":
            max_size = int(config.get("chunk_size", 1200))
            if len(slide_text) > max_size:
                parts = _split_recursive(slide_text, config)
            else:
                parts = [slide_text]
        else:
            parts = _apply_splitter(slide_text, config)

        for part in parts:
            chunks.append({
                "text": part,
                "page": None,
                "slide": idx if bool(config.get("keep_slide_number", True)) else None,
                "chunk_index": chunk_index,
            })
            chunk_index += 1

    if not chunks:
        print(f"⚠️ pptx 無文字內容，略過索引: {row.get('Title')}")

    return [
        {
            "text": c["text"],
            "page": None,
            "slide": c["slide"] if keep_metadata else None,
            "chunk_index": c["chunk_index"],
        }
        for c in chunks
    ]


def _load_chunks(row: Dict[str, Any], config: Dict[str, Any]) -> List[Dict[str, Any]]:
    rel = row.get("FilePath") or ""

    try:
        full_path = _safe_material_path(rel)
    except Exception as e:
        print(f"❌ invalid file path: {rel}, err={e}")
        return []

    if not full_path.exists():
        print(f"❌ missing file: {full_path}")
        return []

    loader = _resolve_loader(row, full_path, config)

    if loader == "pdf":
        return _load_pdf_chunks(full_path, row, config)

    if loader in ("txt", "note", "md"):
        return _load_text_like_chunks(full_path, row, config)

    if loader == "docx":
        return _load_docx_chunks(full_path, row, config)

    if loader == "pptx":
        return _load_pptx_chunks(full_path, row, config)

    print(f"⚠️ 不支援的 loader/type: {loader} (title={row.get('Title')})")
    return []


# ----------------------------
# Index
# ----------------------------
def index_material(material_id: str) -> None:
    """
    針對單一 MaterialId (GUID) 做向量化 + 寫進 Qdrant。
    """
    vs = VectorStore()
    vs.ensure_collection()

    row = _fetch_material(material_id)
    config = _resolve_config(row)
    chunks = _load_chunks(row, config)

    if not chunks:
        print(f"⚠️ Material {material_id} has 0 valid chunks, skip upsert")
        return

    docs = []
    material_id_str = str(row["MaterialId"])

    # loader 先算一次，不要在每個 chunk 都重算
    try:
        loader_name = _resolve_loader(
            row,
            _safe_material_path(row.get("FilePath") or ""),
            config
        )
    except Exception:
        loader_name = str(row.get("Type") or "unknown").lower()

    for ch in chunks:
        text = (ch.get("text") or "").strip()
        if not text:
            continue

        # print(f"\n--- chunk preview material={material_id_str} idx={ch.get('chunk_index')} ---")
        # print(text[:500])
        # print("--- end preview ---\n")

        # ✅ 新增：略過明顯亂碼 / 無效 chunk
        # if is_gibberish(text):
        #     print(
        #         f"⚠️ skip gibberish chunk: material_id={material_id_str}, "
        #         f"chunk_index={ch.get('chunk_index')}"
        #     )
        #     continue

        try:
            emb = embed_text(text)
        except Exception as e:
            print(
                f"❌ embed failed: material_id={material_id_str}, "
                f"chunk_index={ch.get('chunk_index')}, err={e}"
            )
            continue

        doc_id = _make_doc_id(material_id_str, int(ch["chunk_index"]), text)

        docs.append({
            "id": doc_id,
            "embedding": emb,
            "text": text,
            "material_id": material_id_str,
            "title": row.get("Title"),
            "type": row.get("Type"),
            "language": row.get("Language"),
            "style": row.get("Style"),
            "file_path": row.get("FilePath"),
            "bone_id": row.get("BoneId"),
            "bone_small_id": row.get("BoneSmallId"),
            "page": ch.get("page"),
            "slide": ch.get("slide"),
            "chunk_index": ch.get("chunk_index"),
            "tags": [
                f"loader:{loader_name}",
                f"splitter:{config.get('splitter')}",
            ],
        })

    if not docs:
        print(f"⚠️ Material {material_id_str} has no embeddable docs")
        return

    vs.upsert_docs(docs)
    print(f"✅ 材料 {material_id_str} 已寫入向量庫，共 {len(docs)} chunks")
def index_all_materials() -> None:
    with get_connection() as conn:
        cur = conn.cursor()
        cur.execute("SELECT MaterialId FROM agent.TeachingMaterial")
        ids = [str(row[0]) for row in cur.fetchall()]

    for mid in ids:
        try:
            index_material(mid)
        except Exception as e:
            print(f"❌ index failed: MaterialId={mid}, err={e}")


if __name__ == "__main__":
    index_all_materials()